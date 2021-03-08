import os
import platform
import sys
 
package_name = 'python-pptx'
os.system(f'pip install {package_name}')
from pptx import Presentation
def function1(a):
    prs = Presentation(a)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                print(text_frame.text)

print("仅支持PPTX后缀的幻灯片文稿转换\n\n")
b=0
path = os.getcwd()
a=os.listdir(path)
for i in a:
    b+=1
    print(str(b)+":"+i)
c=int(input("请输入所要转换文字的幻灯片文件序号:"))
d=c-1
if(platform.system()=='Windows'):
    print('Windows系统')
    e=path+"\\"+a[d] #windows
    function1(e)
if(platform.system()=='Darwin'):
    print('Mac系统')
    e=a[d] #mac
    function1(e)

print("\n\n\n")

    
while True:
    a= input("输入1退出程序:" )
    if a=="1":
        break
    else:
        print("请勿输入其他字符")
